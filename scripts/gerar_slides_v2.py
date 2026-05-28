#!/usr/bin/env python3
"""
Gera slide deck com imagens explicativas da formula R = K / te.
"""

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np
import math, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─── Constants ───────────────────────────────────────────────────────
BLOCKS_PER_YEAR = 52596
REF_COMMIT = 1_000_000_000
OUTPUT_DIR = "/tmp/ewatts_slides"
os.makedirs(OUTPUT_DIR, exist_ok=True)

# ─── Generate Charts ─────────────────────────────────────────────────

def chart_emission_vs_network():
    """R = K/te curve showing emission vs network size."""
    fig, ax = plt.subplots(figsize=(10, 5.5))
    
    K = 1.9e13
    te_values = np.logspace(7, 13, 100)  # 0.01 miner to 10K miners
    r_per_block = K / te_values
    r_per_year = r_per_block * BLOCKS_PER_YEAR / 1e9  # billions
    
    ax.plot(te_values / REF_COMMIT, r_per_year, 'b-', linewidth=2.5)
    ax.fill_between(te_values / REF_COMMIT, r_per_year, alpha=0.1, color='blue')
    
    # Annotations
    annotations = [
        (0.01, 100, "0.01 miner\n100B Ewatt/ano"),
        (0.1, 10, "0.1 miner\n10B Ewatt/ano"),
        (1, 1, "1 miner\n1B Ewatt/ano"),
        (10, 0.1, "10 miners\n100M Ewatt/ano"),
        (100, 0.01, "100 miners\n10M Ewatt/ano"),
        (1000, 0.001, "1K miners\n1M Ewatt/ano"),
    ]
    for x, y, label in annotations:
        ax.annotate(label, xy=(x, y), fontsize=9,
                   xytext=(x*1.5, y*1.3 if y > 0.001 else y*3),
                   arrowprops=dict(arrowstyle='->', color='gray', lw=1),
                   bbox=dict(boxstyle='round,pad=0.3', facecolor='lightyellow', alpha=0.8))
    
    ax.set_xscale('log')
    ax.set_yscale('log')
    ax.set_xlabel('Tamanho da Rede (miners equivalentes a 1 GB/s)', fontsize=12)
    ax.set_ylabel('Emissao Anual (bilhoes de Ewatt)', fontsize=12)
    ax.set_title('R = K / te  |  K = 1,9 × 10¹³', fontsize=16, fontweight='bold')
    ax.grid(True, alpha=0.3)
    ax.set_xlim(0.005, 5000)
    ax.set_ylim(0.0005, 500)
    
    plt.tight_layout()
    path = os.path.join(OUTPUT_DIR, "emission_vs_network.png")
    fig.savefig(path, dpi=150, bbox_inches='tight')
    plt.close(fig)
    return path

def chart_supply_trajectory():
    """Supply growth over 10 years with network evolution."""
    fig, ax1 = plt.subplots(figsize=(10, 5.5))
    
    K = 1.9e13
    years = np.arange(0, 11)
    supply = 100.0  # genesis
    supplies = [supply]
    emissions = [0]
    miners = [0]
    
    for year in range(1, 11):
        if year <= 1:
            te = int(REF_COMMIT * 0.3)
            m = 0.3
        elif year <= 2:
            te = int(REF_COMMIT * 1)
            m = 1
        elif year <= 3:
            te = int(REF_COMMIT * 3)
            m = 3
        elif year <= 5:
            te = int(REF_COMMIT * 10)
            m = 10
        else:
            te = int(REF_COMMIT * 50)
            m = 50
        
        r_year = (K / te) * BLOCKS_PER_YEAR
        supply += r_year
        supplies.append(supply / 1e9)  # billions
        emissions.append(r_year / 1e9)
        miners.append(m)
    
    # Supply curve
    color1 = '#005B96'
    ax1.fill_between(years, supplies, alpha=0.15, color=color1)
    ax1.plot(years, supplies, 'o-', color=color1, linewidth=2.5, markersize=8, label='Supply (B Ewatt)')
    ax1.set_xlabel('Ano', fontsize=12)
    ax1.set_ylabel('Supply (bilhoes de Ewatt)', fontsize=12, color=color1)
    ax1.tick_params(axis='y', labelcolor=color1)
    ax1.set_xlim(0, 10)
    
    # Annotations on supply
    for i, (y, s) in enumerate(zip(years, supplies)):
        if i in [0, 1, 3, 5, 10]:
            ax1.annotate(f'{s:.2f}B', xy=(y, s), fontsize=9,
                       xytext=(y+0.2, s+max(supplies)*0.02),
                       bbox=dict(boxstyle='round,pad=0.2', facecolor='white', alpha=0.7))
    
    # Emission bars (secondary)
    ax2 = ax1.twinx()
    color2 = '#CC3333'
    ax2.bar(years[1:], emissions[1:], alpha=0.3, color=color2, width=0.6, label='Emissao anual')
    ax2.set_ylabel('Emissao Anual (bilhoes de Ewatt)', fontsize=12, color=color2)
    ax2.tick_params(axis='y', labelcolor=color2)
    
    # Miner count annotations
    for i, (y, m) in enumerate(zip(years[1:], miners[1:])):
        ax1.annotate(f'{m:.0f} miners', xy=(y, supplies[y]), fontsize=8,
                   color='green', fontweight='bold')
    
    ax1.set_title('Trajetoria do Supply — 10 Anos (K = 1,9 × 10¹³)', fontsize=14, fontweight='bold')
    ax1.grid(True, alpha=0.2)
    
    plt.tight_layout()
    path = os.path.join(OUTPUT_DIR, "supply_trajectory.png")
    fig.savefig(path, dpi=150, bbox_inches='tight')
    plt.close(fig)
    return path

def chart_equilibrium_diagram():
    """Diagram showing the market equilibrium cycle."""
    fig, ax = plt.subplots(figsize=(10, 5.5))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 10)
    ax.axis('off')
    
    # Steps in the cycle
    steps = [
        (1, 8, "Preco > Custo\nLucro para miners"),
        (3, 8, "Novos miners\nentram na rede"),
        (5, 8, "te aumenta\n(mais rede)"),
        (7, 8, "R = K/te diminui\n(menos Ewatt/bloco)"),
        (9, 8, "Receita por miner\ncai → equilibrio"),
        (5, 3, "No equilibrio:\nCusto marginal de\nproducao = Preco"),
    ]
    
    for x, y, text in steps:
        if x == 5 and y == 3:
            # Center node
            circle = plt.Circle((x, y), 1.2, color='#005B96', alpha=0.15, ec='#005B96', lw=2)
            ax.add_patch(circle)
            ax.text(x, y, text, ha='center', va='center', fontsize=10, fontweight='bold', color='#005B96')
        else:
            # Step boxes
            rect = plt.Rectangle((x-0.9, y-0.5), 1.8, 1.0, facecolor='#E8F0FE', ec='#005B96', lw=1.5, alpha=0.9)
            ax.add_patch(rect)
            ax.text(x, y, text, ha='center', va='center', fontsize=9, fontweight='bold')
    
    # Arrows
    arrow_props = dict(arrowstyle='->', color='#005B96', lw=2)
    ax.annotate('', xy=(2, 8), xytext=(1.9, 8), arrowprops=arrow_props)
    ax.annotate('', xy=(4, 8), xytext=(3.9, 8), arrowprops=arrow_props)
    ax.annotate('', xy=(6, 8), xytext=(5.9, 8), arrowprops=arrow_props)
    ax.annotate('', xy=(8, 8), xytext=(7.9, 8), arrowprops=arrow_props)
    # Down and back
    ax.annotate('', xy=(9, 4.5), xytext=(9, 6.5), arrowprops=dict(arrowstyle='->', color='red', lw=2))
    ax.annotate('', xy=(6.5, 3), xytext=(8.5, 3.5), arrowprops=dict(arrowstyle='->', color='#005B96', lw=1.5))
    ax.annotate('', xy=(4.5, 6), xytext=(5, 4.5), arrowprops=dict(arrowstyle='->', color='#005B96', lw=1.5))
    # Loop back
    ax.annotate('', xy=(1, 6.5), xytext=(1, 7.5), arrowprops=dict(arrowstyle='->', color='red', lw=2))
    ax.annotate('', xy=(4, 7), xytext=(1.5, 7), arrowprops=dict(arrowstyle='->', color='red', lw=1.5, linestyle='dashed'))
    
    ax.text(5, 1.5, 'Protocolo nao interfere. So K fixo e te variando.', 
           ha='center', fontsize=12, style='italic', color='gray')
    
    path = os.path.join(OUTPUT_DIR, "equilibrium_diagram.png")
    fig.savefig(path, dpi=150, bbox_inches='tight')
    plt.close(fig)
    return path

def chart_k_comparison():
    """Bar chart comparing different K values."""
    fig, ax = plt.subplots(figsize=(10, 5.5))
    
    k_values = ['5e12\n(Conserv.)', '1.9e13\n(Recomend.)', '1e14\n(Agressivo)', '1e15\n(Extremo)']
    k_colors = ['#66B2FF', '#005B96', '#FF8C00', '#CC3333']
    
    # Year 1 emission (1 miner)
    y1 = [263e6, 1e9, 5.26e9, 52.6e9]
    # Year 6 emission (50 miners)
    y6 = [5.26e6, 20e6, 105e6, 1.05e9]
    
    x = np.arange(len(k_values))
    width = 0.35
    
    bars1 = ax.bar(x - width/2, [y/1e9 for y in y1], width, label='Ano 1 (1 miner)', color=k_colors, alpha=0.8)
    bars2 = ax.bar(x + width/2, [y/1e9 for y in y6], width, label='Ano 6 (50 miners)', color=k_colors, alpha=0.4, hatch='//')
    
    ax.set_xticks(x)
    ax.set_xticklabels(k_values, fontsize=11)
    ax.set_ylabel('Emissao (bilhoes de Ewatt/ano)', fontsize=12)
    ax.set_title('Comparacao de Cenarios — K', fontsize=14, fontweight='bold')
    ax.legend(fontsize=11)
    ax.grid(True, alpha=0.2, axis='y')
    
    # Value labels
    for bar in bars1:
        h = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., h, f'{h:.1f}B', 
               ha='center', va='bottom', fontsize=9, fontweight='bold')
    
    plt.tight_layout()
    path = os.path.join(OUTPUT_DIR, "k_comparison.png")
    fig.savefig(path, dpi=150, bbox_inches='tight')
    plt.close(fig)
    return path

def chart_energy_anchor():
    """Energy anchor stability comparison: eWatts vs Bitcoin."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 5))
    
    years = np.arange(0, 11)
    
    # eWatts: DDR latency drift ~1.5%/year
    ewatt_cost = 100 * (1 - 0.015) ** years
    # Bitcoin: ASIC efficiency ~25%/year
    btc_cost = 100 * (1 - 0.25) ** years
    
    ax1.plot(years, ewatt_cost, 'b-', linewidth=2.5, label='eWatts (~1,5%/ano)')
    ax1.fill_between(years, ewatt_cost, alpha=0.15, color='blue')
    ax1.set_title('eWatts: Energy Anchor\n(DDR Latency)', fontsize=13, fontweight='bold')
    ax1.set_xlabel('Anos')
    ax1.set_ylabel('Custo de Producao (%)')
    ax1.set_ylim(0, 105)
    ax1.grid(True, alpha=0.3)
    ax1.legend(fontsize=10)
    ax1.annotate('~14% em 10 anos', xy=(10, ewatt_cost[-1]), fontsize=10,
               xytext=(7, 40), arrowprops=dict(arrowstyle='->', color='blue'))
    
    ax2.plot(years, btc_cost, 'orange', linewidth=2.5, label='Bitcoin (~25%/ano)')
    ax2.fill_between(years, btc_cost, alpha=0.15, color='orange')
    ax2.set_title('Bitcoin: Energy Anchor\n(SHA-256 ASIC)', fontsize=13, fontweight='bold')
    ax2.set_xlabel('Anos')
    ax2.set_ylabel('Custo de Producao (%)')
    ax2.set_ylim(0, 105)
    ax2.grid(True, alpha=0.3)
    ax2.legend(fontsize=10)
    ax2.annotate('~94% em 10 anos', xy=(10, btc_cost[-1]), fontsize=10,
               xytext=(7, 60), arrowprops=dict(arrowstyle='->', color='orange'))
    
    plt.suptitle('Estabilidade da Ancora Energetica', fontsize=14, fontweight='bold', y=1.02)
    plt.tight_layout()
    path = os.path.join(OUTPUT_DIR, "energy_anchor.png")
    fig.savefig(path, dpi=150, bbox_inches='tight')
    plt.close(fig)
    return path

# ─── Generate all charts ─────────────────────────────────────────────
print("Generating charts...")
chart_paths = {
    "emission": chart_emission_vs_network(),
    "supply": chart_supply_trajectory(),
    "equilibrium": chart_equilibrium_diagram(),
    "k_comp": chart_k_comparison(),
    "anchor": chart_energy_anchor(),
}
print("Charts generated.")

# ─── Build PPTX ──────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLACK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x66, 0x66, 0x66)
BLUE = RGBColor(0x00, 0x5B, 0x96)
GREEN = RGBColor(0x00, 0x7A, 0x33)
RED = RGBColor(0xCC, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def title_slide(text, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER

def text_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BLACK
    # Content
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(12), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, is_bold, sz, color, space_after) in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        if sz and sz > 0:
            p.font.size = Pt(sz)
        p.font.bold = is_bold
        if color:
            p.font.color.rgb = color
        p.space_after = Pt(space_after or 6)

def image_slide(title, img_path, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLACK
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY
    
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(0.9), Inches(12.3), Inches(6.0))

# ═══════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════

# Slide 1: Title
title_slide("Formula de Emissao do eWatts", "R = K / te  —  Uma formula, forcas de mercado")

# Slide 2: Formula explanation
text_slide("A Formula", [
    ("R = K / te", True, 28, BLUE, 12),
    ("", False, 0, None, 6),
    ("Onde:", True, 18, BLACK, 6),
    ("  R = Ewatt emitidos por bloco", False, 16, BLACK, 4),
    ("  K = constante fixa do protocolo (definida no genesis)", False, 16, BLACK, 4),
    ("  te = total effective commitment da rede (~miners × 1 GB/s)", False, 16, BLACK, 4),
    ("", False, 0, None, 6),
    ("Propriedades:", True, 18, BLACK, 6),
    ("  Uma unica formula para toda faixa — sem bootstrap vs maduro", False, 16, BLACK, 4),
    ("  Sem percentual do supply, target de inflacao, ou asintota", False, 16, BLACK, 4),
    ("  O mercado regula o equilibrio: mais rede = menos emissao", False, 16, BLACK, 4),
    ("  K e definido uma vez no genesis. Nunca muda.", False, 16, GRAY, 4),
])

# Slide 3: Emission curve chart
image_slide("Emissao vs Tamanho da Rede", chart_paths["emission"],
           "K = 1,9 × 10¹³ — emissao explode com rede pequena, colapsa naturalmente com crescimento")

# Slide 4: K calibration table + K comparison chart
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Calibracao de K"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = BLACK
p2 = tf.add_paragraph()
p2.text = "K define a emissao no boot. O mercado ajusta o resto."
p2.font.size = Pt(14)
p2.font.color.rgb = GRAY

# K comparison chart (left half)
slide4.shapes.add_picture(chart_paths["k_comp"], Inches(0.3), Inches(0.8), Inches(6.5), Inches(3.8))

# Table (right half)
table_shape = slide4.shapes.add_table(6, 5, Inches(7), Inches(0.9), Inches(6), Inches(3.5))
table = table_shape.table
for i, w in enumerate([2.2, 1.5, 1.3, 1.3, 1.3]):
    table.columns[i].width = Inches(w)

headers = ["Cenario", "K", "1 miner", "10 miners", "100 miners"]
data = [
    ["Conservador", "5e12", "263M", "26M", "2,6M"],
    ["Recomendado", "1,9e13", "1,0B", "100M", "10M"],
    ["Agressivo", "1e14", "5,26B", "526M", "52,6M"],
    ["Extremo", "1e15", "52,6B", "5,26B", "526M"],
]

for col_idx, h in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = h
    for para in cell.text_frame.paragraphs:
        para.font.size = Pt(11)
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE

for row_idx, row_data in enumerate(data):
    for col_idx, text in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = text
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(11)
            para.alignment = PP_ALIGN.CENTER
    if row_idx == 1:
        for col_idx in range(5):
            table.cell(row_idx + 1, col_idx).fill.solid()
            table.cell(row_idx + 1, col_idx).fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)

# Note
txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
lines = [
    ("K = 1,9 × 10¹³ e o ponto de partida recomendado.", True, 16, BLUE, 6),
    ("", False, 0, None, 4),
    ("Gera ~$1B/ano de emissao com 1 miner a $1/eWatt.", False, 14, BLACK, 4),
    ("Incentivo real na casa de centenas de milhoes a bilhoes.", False, 14, GRAY, 4),
    ("", False, 0, None, 4),
    ("Referencia: economia global ~$600T. 0,01-0,1% disso = $60M a $600B.", False, 13, GRAY, 4),
]
for i, (text, bld, sz, clr, sp) in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = text
    if sz: p.font.size = Pt(sz)
    p.font.bold = bld
    if clr: p.font.color.rgb = clr
    p.space_after = Pt(sp or 4)

# Slide 5: Supply trajectory
image_slide("Trajetoria do Supply — 10 Anos", chart_paths["supply"],
           "Genesis: 100 Ewatt. Rede cresce de 0,3 a 50 miners. K = 1,9 × 10¹³")

# Slide 6: Market equilibrium diagram
image_slide("Convergencia: Forcas de Mercado", chart_paths["equilibrium"],
           "O protocolo nao interfere — so K fixo e te variando. O mercado acha o equilibrio.")

# Slide 7: Energy anchor
image_slide("Ancora Energetica: eWatts vs Bitcoin", chart_paths["anchor"],
           "DDR latency: ~1,5%/ano de deriva. SHA-256 ASIC: ~25%/ano. eWatts e 10× mais estavel.")

# Slide 8: Summary
text_slide("Resumo", [
    ("R = K / te", True, 28, BLUE, 12),
    ("", False, 0, None, 4),
    ("  - Uma formula, um parametro, sem regimes", False, 18, BLACK, 6),
    ("  - Sem bootstrap vs maduro, sem percentual, sem asintota", False, 18, BLACK, 6),
    ("", False, 0, None, 6),
    ("Mercado regula o equilibrio:", True, 20, BLACK, 8),
    ("  - Mais rede → te sobe → R cai → menos incentivo → equilibrio", False, 16, BLACK, 4),
    ("  - Menos rede → te desce → R sobe → mais incentivo → equilibrio", False, 16, BLACK, 4),
    ("", False, 0, None, 6),
    ("No equilibrio: custo marginal de mineracao ≈ preco do Ewatt.", True, 16, GREEN, 8),
    ("", False, 0, None, 6),
    ("Deriva estrutural de ~1,5%/ano (DDR latency),", False, 14, GRAY, 2),
    ("contra 25-30%/ano do Bitcoin. O anchor e estavel.", False, 14, GRAY, 2),
])

# Save
output_path = "/home/claw/.openclaw/workspace/gustavo_dropbox/eWatts/Analysis/emission_formula_slides.pptx"
prs.save(output_path)
print(f"\nSlide deck salvo: {output_path}")
print(f"Total de slides: {len(prs.slides)}")
print(f"Charts em: {OUTPUT_DIR}")
