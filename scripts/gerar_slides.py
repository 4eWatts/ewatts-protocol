#!/usr/bin/env python3
"""
Gera slide deck com as contas principais da formula de emissao R = K / te.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import math

# ─── Constants ───────────────────────────────────────────────────────
BLOCKS_PER_YEAR = 52596
REF_COMMIT = 1_000_000_000  # 1 GB/s miner em COMMIT_PRECISION

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_slide(title_text, content_lines, left=0.5, top=1.2, width=12.3, height=5.8):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    
    # Content
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, (text, is_bold, size_pt, color) in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        if size_pt > 0:
            if size_pt > 0: p.font.size = Pt(size_pt)
        p.font.bold = is_bold
        if color:
            if color: p.font.color.rgb = color
        p.space_after = Pt(6)
    
    return slide

BLACK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x66, 0x66, 0x66)
BLUE = RGBColor(0x00, 0x5B, 0x96)
GREEN = RGBColor(0x00, 0x7A, 0x33)
RED = RGBColor(0xCC, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1: Formula
# ═══════════════════════════════════════════════════════════════════════
add_slide("Formula de Emissao: R = K / te", [
    ("", False, 0, BLACK),
    ("", False, 0, BLACK),
    ("Ewatt emitidos por bloco = K / te", True, 22, BLUE),
    ("", False, 0, BLACK),
    ("Onde:", True, 18, BLACK),
    ("  R = Ewatt emitidos por bloco", False, 16, BLACK),
    ("  K = constante fixa do protocolo (definida no genesis)", False, 16, BLACK),
    ("  te = total effective commitment da rede (~miners × 1 GB/s)", False, 16, BLACK),
    ("", False, 0, BLACK),
    ("Propriedades:", True, 18, BLACK),
    ("  - Nao ha bootstrap vs maduro: uma unica formula para toda faixa", False, 16, BLACK),
    ("  - Nao ha percentual do supply, target de inflacao, ou asintota", False, 16, BLACK),
    ("  - O mercado regula o equilibrio: mais rede = menos emissao", False, 16, BLACK),
    ("  - O protocolo nao interfere: so K fixo e te variando", False, 16, BLACK),
    ("", False, 0, BLACK),
    ("", False, 0, BLACK),
    ("O K e a UNICA alavanca de politica monetaria do protocolo.", True, 18, BLUE),
    ("Definido uma vez no genesis. Nunca muda.", False, 16, GRAY),
], top=1.3)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2: Calibracao de K
# ═══════════════════════════════════════════════════════════════════════
content = [
    ("Calibracao de K: Emissao Anual por Cenario", True, 28, BLACK),
    ("", False, 0, BLACK),
    ("", False, 0, BLACK),
    ("K define quanto 1 miner de 1 GB/s recebe. O resto escala com 1/te.", False, 16, GRAY),
    ("", False, 0, BLACK),
]

# Table header
columns = [
    ("Cenario", 3.5),
    ("K", 1.5),
    ("1 miner/ano", 1.8),
    ("10 miners/ano", 1.8),
    ("100 miners/ano", 1.8),
    ("1K miners/ano", 1.8),
]
rows_data = [
    ("Boot alto", "5e12", "263M Ew", "26,3M Ew", "2,63M Ew", "263K Ew"),
    ("Boot medio-alto", "1e13", "526M Ew", "52,6M Ew", "5,26M Ew", "526K Ew"),
    ("Boot alto (alvo)", "1,9e13", "1,0B Ew", "100M Ew", "10M Ew", "1,0M Ew"),
    ("Boot muito alto", "1e14", "5,26B Ew", "526M Ew", "52,6M Ew", "5,26M Ew"),
    ("Boot extremo", "1e15", "52,6B Ew", "5,26B Ew", "526M Ew", "52,6M Ew"),
]

# Manual table
x_start = 0.5
y_start = 3.0
col_w = [3.5, 1.5, 1.8, 1.8, 1.8, 1.8]
row_h = 0.45

# Note about K=1e19
content.append(("", False, 0, BLACK))
content.append(("K = 1,9e13 → 1B Ewatt/ano com 1 miner → parece o range certo.", True, 18, BLUE))
content.append(("", False, 0, BLACK))
content.append(("Em $ a $1/eWatt: $1B/ano de emissao no primeiro ano.", False, 16, BLACK))

slide = add_slide("Calibracao de K", content, top=1.0)

# Add table using shapes
from pptx.util import Inches, Pt

# Redo with proper table
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Calibracao de K — Emissao Anual em Ewatt"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "K define quanto 1 miner de 1 GB/s recebe. O resto escala com 1/te."
p.font.size = Pt(14)
p.font.color.rgb = GRAY

# Create table
rows = 6
cols = 6
table_shape = slide2.shapes.add_table(rows, cols, Inches(0.5), Inches(1.6), Inches(12), Inches(3.2))
table = table_shape.table

# Set column widths
for i, w in enumerate([3.5, 1.8, 2.0, 2.0, 2.0, 2.0]):
    table.columns[i].width = Inches(w)

# Header
headers = ["Cenario", "K", "1 miner", "10 miners", "100 miners", "1K miners"]
data_rows = [
    ["Boot alto", "5 × 10¹²", "263M Ewatt", "26,3M Ewatt", "2,63M Ewatt", "263K Ewatt"],
    ["Boot medio-alto", "1 × 10¹³", "526M Ewatt", "52,6M Ewatt", "5,26M Ewatt", "526K Ewatt"],
    ["Boot ideal (alvo)", "1,9 × 10¹³", "1,0B Ewatt", "100M Ewatt", "10M Ewatt", "1,0M Ewatt"],
    ["Boot muito alto", "1 × 10¹⁴", "5,26B Ewatt", "526M Ewatt", "52,6M Ewatt", "5,26M Ewatt"],
    ["Boot extremo", "1 × 10¹⁵", "52,6B Ewatt", "5,26B Ewatt", "526M Ewatt", "52,6M Ewatt"],
]

for col_idx, header in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = header
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.bold = True
        paragraph.font.color.rgb = WHITE
        paragraph.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE

for row_idx, row_data in enumerate(data_rows):
    for col_idx, text in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.alignment = PP_ALIGN.CENTER
            if col_idx == 2:  # highlight target
                paragraph.font.bold = True
                paragraph.font.color.rgb = BLUE
        if row_idx == 2:  # highlight target row
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)

# Note
txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "K = 1,9 × 10¹³ → 1B Ewatt/ano com 1 miner → range ideal."
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = BLUE
p2 = tf.add_paragraph()
p2.text = "Em $ a $1/eWatt: ~$1B/ano de emissao no primeiro ano. Incentivo real para miners."
p2.font.size = Pt(14)
p2.font.color.rgb = GRAY

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3: Emissao vs Rede
# ═══════════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Emissao por Tamanho da Rede (K = 1,9 × 10¹³)"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

# Table
table_shape = slide3.shapes.add_table(8, 5, Inches(0.5), Inches(1.3), Inches(12), Inches(3.8))
table = table_shape.table
for i, w in enumerate([2.5, 2.0, 2.5, 2.5, 2.5]):
    table.columns[i].width = Inches(w)

headers3 = ["Tamanho da Rede", "te", "Ewatt/bloco", "Ewatt/ano", "Equivalente $"]
data3 = [
    ["0,01 miner (boot)", "1 × 10⁷", "1.900.000", "100B", "$100B"],
    ["0,1 miner", "1 × 10⁸", "190.000", "10B", "$10B"],
    ["1 miner = 1 GB/s", "1 × 10⁹", "19.000", "1,0B", "$1,0B"],
    ["3 miners", "3 × 10⁹", "6.333", "333M", "$333M"],
    ["10 miners", "1 × 10¹⁰", "1.900", "100M", "$100M"],
    ["100 miners", "1 × 10¹¹", "190", "10M", "$10M"],
    ["1.000 miners", "1 × 10¹²", "19", "1M", "$1M"],
]

for col_idx, header in enumerate(headers3):
    cell = table.cell(0, col_idx)
    cell.text = header
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.bold = True
        paragraph.font.color.rgb = WHITE
        paragraph.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE

for row_idx, row_data in enumerate(data3):
    for col_idx, text in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.alignment = PP_ALIGN.CENTER
    if row_idx == 2:  # highlight 1 miner row
        for col_idx in range(5):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)

# Note
txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "A emissao cai 10× a cada 10× de rede. Sempre R = K/te, sem condicoes."
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = BLUE
p2 = tf.add_paragraph()
p2.text = "O colapso e natural: se 1.000 miners entram, emissao cai de 1B para 1M Ewatt/ano."
p2.font.size = Pt(14)
p2.font.color.rgb = GRAY

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4: Market Equilibrium
# ═══════════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Convergencia: Forcas de Mercado"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

content4 = [
    ("", False, 0, BLACK),
    ("Ciclo de equilibrio (sem interferencia do protocolo):", True, 20, BLACK),
    ("", False, 0, BLACK),
    ("  1. Preco do Ewatt > custo de mineracao → lucro", False, 16, BLACK),
    ("  2. Lucro atrai novos miners → te aumenta", False, 16, BLACK),
    ("  3. te aumenta → R = K/te diminui", False, 16, BLACK),
    ("  4. R diminui → receita por miner cai → lucro some", False, 16, BLACK),
    ("  5. Entrada de miners para → equilibrio", False, 16, BLACK),
    ("", False, 0, BLACK),
    ("  Se preco < custo: processo inverso. Miners saem. R sobre.", False, 16, BLACK),
    ("", False, 0, BLACK),
    ("", False, 0, BLACK),
    ("A ancoragem energetica e um RESULTADO, nao um target:", True, 20, BLUE),
    ("", False, 0, BLACK),
    ("  No equilibrio, o custo marginal de producao do ultimo miner", False, 16, BLACK),
    ("  que entrou = preco de mercado do Ewatt. Isso e fisica,", False, 16, BLACK),
    ("  nao politica monetaria.", False, 16, BLACK),
    ("", False, 0, BLACK),
    ("", False, 0, BLACK),
    ("Deriva estrutural: DDR latency melhora ~1,5%/ano", True, 18, GRAY),
    ("O custo de producao cai 1,5%/ano — muito menor que", False, 16, GRAY),
    ("os 25-30%/ano de Bitcoin (ASIC efficiency).", False, 16, GRAY),
]

for i, (text, is_bold, size_pt, color) in enumerate(content4):
    if i == 0:
        txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    if size_pt > 0: p.font.size = Pt(size_pt)
    p.font.bold = is_bold
    if color: p.font.color.rgb = color
    p.space_after = Pt(4)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5: Trajetoria do Supply (10 anos)
# ═══════════════════════════════════════════════════════════════════════
# Simulate supply trajectory
def simulate(K, scenario_years):
    """Returns list of (year, supply, emission, miners)"""
    supply = 100.0  # genesis Ewatt
    results = [(0, supply, 0, 0)]
    
    for year in range(1, scenario_years + 1):
        # Network growth scenario
        if year <= 1:
            te = int(REF_COMMIT * 0.3)  # 0.3 miner effective in year 1
            miners = 0.3
        elif year <= 2:
            te = int(REF_COMMIT * 1)    # 1 miner
            miners = 1
        elif year <= 3:
            te = int(REF_COMMIT * 3)    # 3 miners
            miners = 3
        elif year <= 5:
            te = int(REF_COMMIT * 10)   # 10 miners
            miners = 10
        else:
            te = int(REF_COMMIT * 50)   # 50 miners
            miners = 50
        
        r_per_block = K / te
        r_per_year = r_per_block * BLOCKS_PER_YEAR
        supply += r_per_year
        
        em_b = r_per_block
        em_a = r_per_year
        
        results.append((year, supply, em_a, miners))
    
    return results

slide5 = prs.slides.add_slide(prs.slide_layouts[6])
txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Trajetoria do Supply em 10 Anos (K = 1,9 × 10¹³)"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

# Table
table_shape = slide5.shapes.add_table(12, 6, Inches(0.5), Inches(1.2), Inches(12), Inches(4.8))
table = table_shape.table
for i, w in enumerate([1.0, 2.0, 2.0, 2.5, 2.0, 2.5]):
    table.columns[i].width = Inches(w)

headers5 = ["Ano", "te", "Miners", "Ewatt/ano", "Supply (Ew)", "Crescimento %"]
data5 = [
    ["0", "0", "0", "0", "100", "—"],
]

sim = simulate(1.9e13, 10)
for year, supply, em_annual, miners in sim[1:]:
    growth = em_annual / (supply - em_annual) * 100 if supply > em_annual else 0
    te_val = int(REF_COMMIT * max(0.3, miners))
    data5.append([
        str(year),
        f"{te_val:.0e}",
        f"{miners:.1f}",
        f"{em_annual/1e6:.1f}M" if em_annual < 1e9 else f"{em_annual/1e9:.1f}B",
        f"{supply/1e6:.1f}M" if supply < 1e9 else f"{supply/1e9:.2f}B",
        f"{growth:.1f}%",
    ])

for col_idx, header in enumerate(headers5):
    cell = table.cell(0, col_idx)
    cell.text = header
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(13)
        paragraph.font.bold = True
        paragraph.font.color.rgb = WHITE
        paragraph.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE

for row_idx, row_data in enumerate(data5):
    for col_idx, text in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.alignment = PP_ALIGN.CENTER

# Note
txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12.3), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Genesis de 100 Ewatt com K = 1,9 × 10¹³. Emissao explode no boot, estabiliza com a rede."
p.font.size = Pt(14)
p.font.color.rgb = GRAY

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6: Comparacao de K
# ═══════════════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
txBox = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Comparacao de Cenarios: Escolha de K"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

table_shape = slide6.shapes.add_table(6, 5, Inches(0.5), Inches(1.3), Inches(12), Inches(3.0))
table = table_shape.table
for i, w in enumerate([3.0, 2.5, 2.5, 2.5, 2.5]):
    table.columns[i].width = Inches(w)

headers6 = ["Cenario (K)", "1 miner (ano 1)", "10 miners (ano 3)", "50 miners (ano 6)", "Supply final (10a)"]
data6 = [
    ["Conservador (5e12)", "263M Ewatt", "26,3M Ewatt", "5,26M Ewatt", "~332M Ewatt"],
    ["Recomendado (1,9e13)", "1,0B Ewatt", "100M Ewatt", "20M Ewatt", "~1,26B Ewatt"],
    ["Agressivo (1e14)", "5,26B Ewatt", "526M Ewatt", "105M Ewatt", "~6,6B Ewatt"],
    ["Extremo (1e15)", "52,6B Ewatt", "5,26B Ewatt", "1,05B Ewatt", "~66B Ewatt"],
]

for col_idx, header in enumerate(headers6):
    cell = table.cell(0, col_idx)
    cell.text = header
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(13)
        paragraph.font.bold = True
        paragraph.font.color.rgb = WHITE
        paragraph.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE

for row_idx, row_data in enumerate(data6):
    for col_idx, text in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.alignment = PP_ALIGN.CENTER
    if row_idx == 1:  # highlight recommended
        for col_idx in range(5):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)

txBox = slide6.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True

lines = [
    ("Referencia de valor: $600T de economia global", True, 18, BLUE),
    ("", False, 0, BLACK),
    ("Se eWatts capturar 0,01-0,1% da economia global no primeiro ano:", False, 16, BLACK),
    ("  → $60M a $600B de market cap potencial", False, 16, BLACK),
    ("  → K = 1,9e13 da $1B/ano de emissao a $1/eWatt", False, 16, BLACK),
    ("  → Range realista: $100M-$10B de emissao no primeiro ano", False, 16, BLACK),
    ("", False, 0, BLACK),
    ("K = 1,9 × 10¹³ parece o ponto de partida ideal.", True, 18, BLUE),
]

for i, (text, is_bold, size_pt, color) in enumerate(lines):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    if size_pt > 0: p.font.size = Pt(size_pt)
    p.font.bold = is_bold
    if color: p.font.color.rgb = color
    p.space_after = Pt(4)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7: Conclusao
# ═══════════════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
txBox = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Resumo"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

content7 = [
    ("", False, 0, BLACK),
    ("Formula: R = K / te", True, 24, BLUE),
    ("", False, 0, BLACK),
    ("  - Um parametro so: K", False, 18, BLACK),
    ("  - Uma formula so para toda faixa da rede", False, 18, BLACK),
    ("  - Sem bootstrap vs maduro, sem percentual, sem asintota", False, 18, BLACK),
    ("", False, 0, BLACK),
    ("  - K = 1,9 × 10¹³ → 1B Ewatt/ano com 1 miner", True, 18, BLUE),
    ("", False, 0, BLACK),
    ("Mercado regula o equilibrio:", True, 20, BLACK),
    ("", False, 0, BLACK),
    ("  - Mais rede (+) → te (+) → R (-) → menos incentivo", False, 18, BLACK),
    ("  - Menos rede (-) → te (-) → R (+) → mais incentivo", False, 18, BLACK),
    ("", False, 0, BLACK),
    ("  - No equilibrio, custo de mineracao ≈ preco do Ewatt", True, 18, GREEN),
    ("", False, 0, BLACK),
    ("Deriva estrutural de ~1,5%/ano (DDR latency),", False, 16, GRAY),
    ("contra 25-30%/ano do Bitcoin. O anchor e estavel.", False, 16, GRAY),
]

for i, (text, is_bold, size_pt, color) in enumerate(content7):
    if i == 0:
        txBox = slide7.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    if size_pt > 0: p.font.size = Pt(size_pt)
    p.font.bold = is_bold
    if color: p.font.color.rgb = color
    p.space_after = Pt(4)

# Save
output_path = "/home/claw/.openclaw/workspace/gustavo_dropbox/eWatts/Analysis/emission_formula_slides.pptx"
prs.save(output_path)
print(f"Slide deck salvo em: {output_path}")
print(f"Total de slides: {len(prs.slides)}")
